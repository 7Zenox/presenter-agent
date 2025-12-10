"""PowerPoint file parser to extract slides."""
import base64
import io
from typing import List
from pptx import Presentation
from app.models import Slide


def parse_powerpoint(file_data: bytes) -> List[Slide]:
    """
    Parse a PowerPoint file and extract slides.
    
    Args:
        file_data: Raw bytes of the PowerPoint file (.ppt or .pptx)
    
    Returns:
        List of Slide objects extracted from the presentation
    """
    try:
        # Create a file-like object from bytes
        ppt_file = io.BytesIO(file_data)
        
        # Load the presentation
        prs = Presentation(ppt_file)
        
        slides = []
        for idx, slide in enumerate(prs.slides, start=1):
            # Extract title (usually from title shape)
            title = ""
            bullets = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    # Check if it's a title (usually first shape or shape with title placeholder)
                    if shape.is_placeholder:
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            title = text
                            continue
                    
                    # If no title found yet and this is the first text shape, use it as title
                    if not title and not bullets:
                        title = text
                    else:
                        # Extract bullet points
                        # Check if shape has paragraphs (for bullet lists)
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    bullets.append(para_text)
                        else:
                            # Simple text shape
                            if text and text != title:
                                bullets.append(text)
            
            # If no title found, use a default
            if not title:
                title = f"Slide {idx}"
            
            # Create slide object
            slide_obj = Slide(
                id=idx,
                title=title,
                bullets=bullets if bullets else []
            )
            slides.append(slide_obj)
        
        return slides
    
    except Exception as e:
        raise Exception(f"Failed to parse PowerPoint file: {e}")


def parse_powerpoint_from_base64(base64_data: str) -> List[Slide]:
    """
    Parse a PowerPoint file from base64 encoded string.
    
    Args:
        base64_data: Base64 encoded PowerPoint file data
    
    Returns:
        List of Slide objects extracted from the presentation
    """
    try:
        # Decode base64 to bytes
        file_data = base64.b64decode(base64_data)
        return parse_powerpoint(file_data)
    except Exception as e:
        raise Exception(f"Failed to decode PowerPoint file: {e}")



