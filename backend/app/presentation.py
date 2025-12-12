"""
PowerPoint presentation parser and tool handlers for the presenter agent.
"""
import json
from typing import Dict, List, Optional
from pptx import Presentation
import io
import logging

logger = logging.getLogger(__name__)


class PresentationManager:
    """Manages PowerPoint presentation state and operations."""
    
    def __init__(self):
        self.presentation: Optional[Presentation] = None
        self.slides: List[Dict] = []
        self.current_slide_index: int = 0
    
    def load_presentation(self, pptx_bytes: bytes) -> Dict:
        """Load a PowerPoint presentation from bytes."""
        try:
            # Clear any existing presentation data first
            logger.info("Clearing any existing presentation data...")
            self.presentation = None
            self.slides = []
            self.current_slide_index = 0
            
            # Load new presentation
            self.presentation = Presentation(io.BytesIO(pptx_bytes))
            self.slides = []
            
            for idx, slide in enumerate(self.presentation.slides):
                slide_data = {
                    "index": idx,
                    "title": self._extract_title(slide, idx),
                    "content": self._extract_content(slide),
                    "notes": self._extract_notes(slide),
                }
                self.slides.append(slide_data)
            
            self.current_slide_index = 0
            
            # Log first slide to verify content
            if self.slides:
                first_slide = self.slides[0]
                logger.info(f"✅ Loaded NEW presentation with {len(self.slides)} slides")
                logger.info(f"   First slide title: '{first_slide['title']}'")
                logger.info(f"   First slide content (first 200 chars): '{first_slide['content'][:200]}'")
                
                # Verify storage: Log a few slides to confirm they're stored
                logger.info(f"   📦 Storage verification:")
                logger.info(f"      - Total slides stored in memory: {len(self.slides)}")
                if len(self.slides) > 4:
                    slide_5 = self.slides[4]  # Index 4 = slide 5
                    logger.info(f"      - Slide 5 stored: ✅ (title: '{slide_5['title'][:50]}...')")
                    logger.info(f"      - Slide 5 content length: {len(slide_5['content'])} chars")
                logger.info(f"   ✅ All slides stored in memory - get_slide() tool can fetch any slide on demand")
            else:
                logger.warning("⚠️ Loaded presentation but no slides found!")
            
            return {
                "success": True,
                "total_slides": len(self.slides),
                "slides": self.slides,
            }
        except Exception as e:
            logger.error(f"Error loading presentation: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return {
                "success": False,
                "error": str(e),
            }
    
    def _extract_title(self, slide, slide_index: int) -> str:
        """Extract title from slide."""
        if slide.shapes.title and slide.shapes.title.text:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
        # Fallback to slide number (1-based) if no title found
        return f"Slide {slide_index + 1}"
    
    def _extract_content(self, slide) -> str:
        """Extract all text content from slide, including bullet points and nested text."""
        content_parts = []
        title_text = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
        
        def extract_text_from_shape(shape):
            """Recursively extract text from a shape and its children."""
            texts = []
            
            # Handle text frames first (for bullet points and paragraphs)
            # This gives us structured content with proper paragraph separation
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text and para_text != title_text:
                        texts.append(para_text)
            # Fallback to direct text only if no text_frame exists
            elif hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text != title_text:
                    texts.append(text)
            
            # Handle grouped shapes (recursively)
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    texts.extend(extract_text_from_shape(sub_shape))
            
            return texts
        
        # Extract text from all shapes
        for shape in slide.shapes:
            # Skip the title shape (we already have it separately)
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            
            shape_texts = extract_text_from_shape(shape)
            content_parts.extend(shape_texts)
        
        # Remove duplicates while preserving order
        seen = set()
        unique_parts = []
        for part in content_parts:
            if part and part not in seen:
                seen.add(part)
                unique_parts.append(part)
        
        return "\n".join(unique_parts)
    
    def _extract_notes(self, slide) -> str:
        """Extract speaker notes from slide."""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""
    
    def get_current_slide(self) -> Dict:
        """Get current slide information."""
        if not self.slides or self.current_slide_index >= len(self.slides):
            return {"error": "No slides available"}
        
        slide = self.slides[self.current_slide_index].copy()
        slide["is_current"] = True
        return slide
    
    def navigate_to_slide(self, action: str, slide_index: Optional[int] = None) -> Dict:
        """Navigate to a slide.
        
        Args:
            action: "next", "prev", or "jump"
            slide_index: Required if action is "jump"
        """
        if not self.slides:
            return {"error": "No presentation loaded"}
        
        total_slides = len(self.slides)
        
        if action == "next":
            if self.current_slide_index < total_slides - 1:
                self.current_slide_index += 1
            else:
                return {"error": "Already on last slide"}
        
        elif action == "prev":
            if self.current_slide_index > 0:
                self.current_slide_index -= 1
            else:
                return {"error": "Already on first slide"}
        
        elif action == "jump":
            if slide_index is None:
                return {"error": "slide_index required for jump action"}
            if 0 <= slide_index < total_slides:
                self.current_slide_index = slide_index
            else:
                return {"error": f"Invalid slide index. Must be between 0 and {total_slides - 1}"}
        
        else:
            return {"error": f"Invalid action: {action}. Use 'next', 'prev', or 'jump'"}
        
        return {
            "success": True,
            "current_slide": self.current_slide_index,
            "total_slides": total_slides,
            "slide": self.get_current_slide(),
        }
    
    def get_slide_content(self, slide_index: int) -> Dict:
        """Get content of a specific slide from memory storage.
        
        This fetches the full slide data (title, content, notes) that was stored
        when the presentation was loaded. All slides are kept in memory for fast access.
        """
        if not self.slides:
            logger.error("❌ get_slide_content called but no slides in memory!")
            return {"error": "No presentation loaded"}
        
        if 0 <= slide_index < len(self.slides):
            slide = self.slides[slide_index].copy()
            slide["is_current"] = (slide_index == self.current_slide_index)
            logger.info(f"   📦 Retrieved slide {slide_index + 1} from memory storage")
            logger.info(f"      Title: '{slide.get('title', 'N/A')}'")
            logger.info(f"      Content length: {len(slide.get('content', ''))} chars")
            return slide
        else:
            logger.error(f"   ❌ Invalid slide_index {slide_index} (valid range: 0-{len(self.slides) - 1})")
            return {"error": f"Invalid slide index. Must be between 0 and {len(self.slides) - 1}"}
    
    def get_all_slides_summary(self) -> str:
        """Get slide numbers, titles, and content previews for AI context.
        
        Returns slide numbers, titles, and a content preview (first 200 chars)
        to give the AI enough context to search and navigate effectively
        while staying within token limits.
        """
        if not self.slides:
            return "No presentation loaded."
        
        import json
        # Create list with slide number, title, and content preview
        slides_list = []
        for slide in self.slides:
            content = slide.get('content', '')
            # Create a content preview - first 200 chars, clean up whitespace
            content_preview = ' '.join(content.split())[:200]
            if len(content) > 200:
                content_preview += "..."
            
            slides_list.append({
                "slide_number": slide['index'] + 1,  # 1-based
                "title": slide.get('title', 'Untitled'),
                "content_preview": content_preview,
            })
        
        slides_json = json.dumps(slides_list, indent=2, ensure_ascii=False)
        
        # Log summary stats
        logger.info(f"📄 Generated slide list: {len(self.slides)} slides (numbers + titles only)")
        logger.info(f"   Summary length: {len(slides_json)} chars (~{len(slides_json) // 4} tokens)")
        
        # Verify all slides are included
        logger.info(f"   🔍 Verifying all slides included:")
        logger.info(f"      Expected slides: {len(self.slides)}")
        logger.info(f"      Actual entries: {len(slides_list)}")
        if len(slides_list) != len(self.slides):
            logger.error(f"   ❌❌❌ MISMATCH: Expected {len(self.slides)} slides but got {len(slides_list)} entries!")
        
        # Log slide numbers to verify they're sequential
        slide_numbers = [s['slide_number'] for s in slides_list]
        logger.info(f"   📋 Slide numbers: {slide_numbers[:10]}{'...' if len(slide_numbers) > 10 else ''}")
        if len(slide_numbers) > 0:
            logger.info(f"   📋 First slide: {slide_numbers[0]}, Last slide: {slide_numbers[-1]}")
            if 5 in slide_numbers:
                slide_5 = next(s for s in slides_list if s['slide_number'] == 5)
                logger.info(f"   ✅ Slide 5 found: title='{slide_5.get('title', 'N/A')}'")
            else:
                logger.error(f"   ❌❌❌ Slide 5 NOT found! Available slides: {slide_numbers}")
        
        return f"""Total slides: {len(self.slides)}

SLIDES (use get_slide(slide_number=X) tool to retrieve full content):
{slides_json}"""
    
    def get_all_slides_json(self) -> str:
        """Get slides as JSON string for embedding in instructions."""
        import json
        if not self.slides:
            return "[]"
        # Add slide_number (1-based) to each slide for clarity
        slides_with_number = []
        for slide in self.slides:
            slide_copy = slide.copy()
            slide_copy['slide_number'] = slide['index'] + 1  # 1-based slide number
            slides_with_number.append(slide_copy)
        return json.dumps(slides_with_number, indent=2)
    
    def get_full_presentation_for_ai(self) -> str:
        """Get ALL slide content formatted for AI instructions.
        
        Returns full content of every slide in a readable text format
        so the AI has all content in context without needing tool calls.
        """
        if not self.slides:
            return "No presentation loaded."
        
        lines = [f"PRESENTATION: {len(self.slides)} slides total\n"]
        lines.append("=" * 50 + "\n")
        
        for slide in self.slides:
            slide_num = slide['index'] + 1
            title = slide.get('title', 'Untitled')
            content = slide.get('content', '').strip()
            notes = slide.get('notes', '').strip()
            
            lines.append(f"--- SLIDE {slide_num}: {title} ---\n")
            if content:
                lines.append(f"CONTENT:\n{content}\n")
            if notes:
                lines.append(f"NOTES:\n{notes}\n")
            lines.append("")
        
        result = "\n".join(lines)
        logger.info(f"📄 Full presentation content: {len(result)} chars (~{len(result) // 4} tokens)")
        return result


# Global presentation manager instance
presentation_manager = PresentationManager()



